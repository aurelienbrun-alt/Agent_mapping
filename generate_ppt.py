"""Generate a professional PowerPoint presentation for the Regulatory Framework Mapper agent."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Wavestone brand palette ──────────────────────────────────────────────────
WS_ORANGE   = RGBColor(0xE6, 0x33, 0x12)   # #E63312  primary
WS_DARK     = RGBColor(0x1A, 0x1A, 0x2E)   # #1A1A2E  near-black navy
WS_NAVY     = RGBColor(0x1E, 0x3A, 0x5F)   # #1E3A5F  dark blue
WS_LIGHT    = RGBColor(0xF5, 0xF5, 0xF5)   # #F5F5F5  off-white bg
WS_MID      = RGBColor(0x4A, 0x4A, 0x6A)   # #4A4A6A  mid-grey text
WS_ACCENT   = RGBColor(0xFF, 0x8C, 0x42)   # #FF8C42  warm orange accent
WS_GREEN    = RGBColor(0x27, 0xAE, 0x60)   # coverage full
WS_YELLOW   = RGBColor(0xF3, 0x9C, 0x12)   # coverage partial
WS_RED      = RGBColor(0xC0, 0x39, 0x2B)   # gap
WS_BLUE_LT  = RGBColor(0x2E, 0x86, 0xAB)   # #2E86AB  azure blue
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE, x, y, w, h)
    # MSO_SHAPE_TYPE is not right for add_shape; use 1 (rectangle)
    return shape


def rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(1.5), radius=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    sp = slide.shapes.add_shape(1, x, y, w, h)   # 1 = MSO_SHAPE.RECTANGLE
    fill = sp.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = sp.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_w
    else:
        line.fill.background()
    return sp


def rounded_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(1.5)):
    sp = slide.shapes.add_shape(5, x, y, w, h)   # 5 = MSO_SHAPE.ROUNDED_RECTANGLE
    fill = sp.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = sp.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_w
    else:
        line.fill.background()
    return sp


def txbox(slide, text, x, y, w, h,
          font_size=Pt(11), bold=False, color=BLACK, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_label(slide, text, x, y, w, h,
              bg=None, txt_color=WHITE, font_size=Pt(10), bold=True,
              align=PP_ALIGN.CENTER, rounded=False):
    if rounded:
        sp = rounded_rect(slide, x, y, w, h, fill_rgb=bg, line_rgb=None)
    else:
        sp = rect(slide, x, y, w, h, fill_rgb=bg, line_rgb=None)
    tf = sp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = txt_color
    return sp


def arrow_right(slide, x, y, length, color=WS_MID, width=Pt(1.5)):
    x, y, length = int(x), int(y), int(length)
    connector = slide.shapes.add_connector(1, x, y, x + length, y)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def arrow_down(slide, x, y, length, color=WS_MID, width=Pt(1.5)):
    x, y, length = int(x), int(y), int(length)
    connector = slide.shapes.add_connector(1, x, y, x, y + length)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


# ── slide builders ───────────────────────────────────────────────────────────

def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    # Full background
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_DARK)

    # Orange accent stripe left
    rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill_rgb=WS_ORANGE)

    # Decorative block top-right
    rect(slide, Inches(10.5), Inches(0), Inches(2.83), Inches(2.5), fill_rgb=WS_NAVY)

    # Orange dot grid (decorative)
    for i in range(3):
        for j in range(4):
            rect(slide, Inches(10.7 + i * 0.55), Inches(0.3 + j * 0.45),
                 Inches(0.18), Inches(0.18), fill_rgb=WS_ORANGE)

    # Title
    txbox(slide, "Regulatory Framework Mapper",
          Inches(0.7), Inches(2.0), Inches(9.5), Inches(1.2),
          font_size=Pt(40), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    txbox(slide, "Architecture & Workflow de l'Agent",
          Inches(0.7), Inches(3.25), Inches(9), Inches(0.7),
          font_size=Pt(22), bold=False, color=WS_ACCENT, align=PP_ALIGN.LEFT)

    # Description
    txbox(slide,
          "Cartographie automatique de référentiels réglementaires\n"
          "via Azure OpenAI · Embeddings sémantiques · Analyse multi-dimensionnelle",
          Inches(0.7), Inches(4.0), Inches(9), Inches(1.0),
          font_size=Pt(14), color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.LEFT)

    # Bottom bar
    rect(slide, Inches(0), Inches(6.9), SLIDE_W, Inches(0.6), fill_rgb=WS_ORANGE)
    txbox(slide, "Wavestone  |  2026  |  Confidentiel",
          Inches(0.5), Inches(6.95), Inches(6), Inches(0.4),
          font_size=Pt(10), color=WHITE, align=PP_ALIGN.LEFT)


def make_overview_slide(prs):
    """Global pipeline — all 6 stages in one view."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_LIGHT)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), fill_rgb=WS_DARK)
    rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill_rgb=WS_ORANGE)

    txbox(slide, "Vue d'ensemble — Pipeline complet",
          Inches(0.4), Inches(0.12), Inches(10), Inches(0.5),
          font_size=Pt(20), bold=True, color=WHITE)

    stages = [
        ("01", "Prétraitement",         "Lecture Excel\nHarmonisation ENISA\nAtomisation & Champs\nEmbeddings",             WS_NAVY),
        ("02", "Mapping directionnel",  "Scoring multi-dim.\nFusion RRF\nÉvaluation LLM\nCache décisions",                  WS_BLUE_LT),
        ("03", "Juge final",            "Révision LLM batch\n/catégorie ENISA\nCorrection FP/FN",                           WS_ORANGE),
        ("04", "Synthèse parent",       "Agrégation atomes\nClassification lacune\nTexte synthèse LLM",                     RGBColor(0x8E, 0x44, 0xAD)),
        ("05", "Workbook Excel",        "8 onglets\nDashboard KPIs\nPiste d'évidence",                                       WS_GREEN),
        ("06", "Rapport analyse",       "Logs JSONL\nStats & timing\nMarkdown report",                                       WS_MID),
    ]

    box_w = Inches(1.9)
    box_h = Inches(4.8)
    gap = Inches(0.22)
    start_x = Inches(0.45)
    start_y = Inches(1.0)

    for i, (num, title, body, color) in enumerate(stages):
        x = start_x + i * (box_w + gap)

        # card background
        rounded_rect(slide, x, start_y, box_w, box_h,
                     fill_rgb=WHITE, line_rgb=color)

        # colored header
        rounded_rect(slide, x, start_y, box_w, Inches(0.85),
                     fill_rgb=color, line_rgb=None)

        # number badge
        add_label(slide, num, x + Inches(0.07), start_y + Inches(0.1),
                  Inches(0.42), Inches(0.42),
                  bg=WHITE, txt_color=color, font_size=Pt(14), bold=True,
                  align=PP_ALIGN.CENTER, rounded=True)

        # stage title
        txbox(slide, title,
              x + Inches(0.54), start_y + Inches(0.12),
              box_w - Inches(0.6), Inches(0.65),
              font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

        # body text
        txbox(slide, body,
              x + Inches(0.12), start_y + Inches(1.0),
              box_w - Inches(0.18), Inches(3.6),
              font_size=Pt(9.5), color=WS_MID, align=PP_ALIGN.LEFT)

        # arrow between cards
        if i < len(stages) - 1:
            ax = x + box_w + Inches(0.03)
            ay = start_y + box_h / 2
            arrow_right(slide, ax, ay, gap - Inches(0.06),
                        color=WS_MID, width=Pt(2))

    # inputs strip bottom
    rect(slide, Inches(0.25), Inches(6.1), Inches(5.5), Inches(0.45),
         fill_rgb=RGBColor(0xEB, 0xF5, 0xFB), line_rgb=WS_BLUE_LT)
    txbox(slide, "ENTRÉES  :  Framework A (xlsx)   ·   Framework B (xlsx)   ·   Guidelines RAG (pdf/txt/md)",
          Inches(0.35), Inches(6.12), Inches(5.3), Inches(0.38),
          font_size=Pt(8.5), color=WS_NAVY, bold=False)

    rect(slide, Inches(7.5), Inches(6.1), Inches(5.6), Inches(0.45),
         fill_rgb=RGBColor(0xEB, 0xF9, 0xF1), line_rgb=WS_GREEN)
    txbox(slide, "SORTIES  :  mapping_{FrA}_{FrB}_{date}.xlsx   ·   run_{...}.jsonl   ·   log_analysis.md",
          Inches(7.6), Inches(6.12), Inches(5.4), Inches(0.38),
          font_size=Pt(8.5), color=RGBColor(0x1A, 0x5C, 0x32), bold=False)

    # bottom bar
    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), fill_rgb=WS_DARK)
    txbox(slide, "Regulatory Framework Mapper  |  Architecture Overview",
          Inches(0.4), Inches(7.0), Inches(8), Inches(0.4),
          font_size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xAA))


def make_preprocessing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_LIGHT)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), fill_rgb=WS_NAVY)
    rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill_rgb=WS_ORANGE)

    add_label(slide, "01", Inches(0.35), Inches(0.13), Inches(0.48), Inches(0.48),
              bg=WS_ORANGE, font_size=Pt(15), bold=True, align=PP_ALIGN.CENTER, rounded=True)
    txbox(slide, "Étape 1 — Prétraitement des référentiels",
          Inches(1.0), Inches(0.12), Inches(10), Inches(0.5),
          font_size=Pt(20), bold=True, color=WHITE)

    # Sub-steps
    steps = [
        ("1a", "Lecture Excel",
         "Lecture de l'onglet Excel configuré\nMapping colonnes : ID · Title · Requirement · Category\nLimite MAX_REQUIREMENTS_PER_FRAMEWORK",
         WS_BLUE_LT),
        ("1b", "Harmonisation ENISA",
         "Normalisation vers les 13 catégories ENISA\n① Règles déterministes (category_overrides.xlsx)\n② LLM en fallback (PROMPT_CATEGORY_HARMONIZATION)\nSeuils : strong ≥ 0.85 · medium ≥ 0.60",
         WS_NAVY),
        ("1c", "Atomisation",
         "LLM (PROMPT_ATOMIZE) : découpe chaque exigence\ncomplexe en obligations unitaires (atomes)\nEx. : 1 exigence → [A1] accès · [A2] audit · [A3] formation",
         WS_ORANGE),
        ("1d", "Extraction de champs",
         "LLM (PROMPT_EXTRACT_FIELDS) par atome :\nactor · action · object · condition · deadline\nevidence · obligation_type · control_type · keywords",
         RGBColor(0x8E, 0x44, 0xAD)),
        ("1e", "Embeddings",
         "Azure OpenAI : text-embedding-3-large\n512 dimensions par atome\nBatches de 64 · Cosinus pour similarité",
         WS_GREEN),
    ]

    bw = Inches(2.35)
    bh = Inches(4.5)
    gap = Inches(0.15)
    sx = Inches(0.4)
    sy = Inches(1.0)

    for i, (num, title, body, color) in enumerate(steps):
        x = sx + i * (bw + gap)
        rounded_rect(slide, x, sy, bw, bh, fill_rgb=WHITE, line_rgb=color)
        rounded_rect(slide, x, sy, bw, Inches(0.7), fill_rgb=color, line_rgb=None)

        add_label(slide, num, x + Inches(0.06), sy + Inches(0.1),
                  Inches(0.42), Inches(0.42), bg=WHITE, txt_color=color,
                  font_size=Pt(12), bold=True, align=PP_ALIGN.CENTER, rounded=True)
        txbox(slide, title, x + Inches(0.54), sy + Inches(0.12),
              bw - Inches(0.6), Inches(0.58),
              font_size=Pt(10.5), bold=True, color=WHITE)
        txbox(slide, body, x + Inches(0.1), sy + Inches(0.85),
              bw - Inches(0.18), Inches(3.5),
              font_size=Pt(9), color=WS_MID)

        if i < len(steps) - 1:
            ax = x + bw + Inches(0.02)
            ay = sy + bh / 2
            arrow_right(slide, ax, ay, gap, color=WS_MID, width=Pt(2))

    # Cache note
    rect(slide, Inches(0.4), Inches(5.75), Inches(12.5), Inches(0.75),
         fill_rgb=RGBColor(0xFF, 0xF8, 0xE1), line_rgb=WS_YELLOW)
    txbox(slide, "💾  Cache  :  docs/cache/{framework_key}/  →  category_harmonization.json  ·  atomized_requirements.json  ·  atomic_requirements_with_fields.json  ·  atomic_requirements.json\n"
                 "Validation SHA256 du fichier source (STRICT_INPUT_CACHE_VALIDATION=true) — prévient la réutilisation silencieuse d'un cache obsolète",
          Inches(0.55), Inches(5.78), Inches(12.2), Inches(0.65),
          font_size=Pt(8.5), color=RGBColor(0x7D, 0x6A, 0x00))

    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), fill_rgb=WS_DARK)
    txbox(slide, "Regulatory Framework Mapper  |  Étape 1 : Prétraitement",
          Inches(0.4), Inches(7.0), Inches(8), Inches(0.4),
          font_size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xAA))


def make_enisa_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_LIGHT)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), fill_rgb=WS_NAVY)
    rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill_rgb=WS_ORANGE)

    txbox(slide, "Classification ENISA — Taxonomie Pivot Cybersécurité",
          Inches(0.35), Inches(0.12), Inches(12.5), Inches(0.5),
          font_size=Pt(20), bold=True, color=WHITE)

    # ── LEFT: 13 categories ───────────────────────────────────────────────
    rect(slide, Inches(0.35), Inches(0.85), Inches(6.15), Inches(6.05),
         fill_rgb=WHITE, line_rgb=WS_NAVY)
    txbox(slide, "13 Catégories Cybersécurité (Taxonomie ENISA / NIS2)",
          Inches(0.45), Inches(0.90), Inches(5.95), Inches(0.35),
          font_size=Pt(11), bold=True, color=WS_NAVY)

    cats = [
        ("01", "Governance & Risk Management"),
        ("02", "Legal & Compliance"),
        ("03", "Asset Management"),
        ("04", "Identity & Access Management"),
        ("05", "Cryptography & Key Management"),
        ("06", "Physical & Environmental Security"),
        ("07", "Network Security"),
        ("08", "Application Security"),
        ("09", "Incident Management"),
        ("10", "Business Continuity & DR"),
        ("11", "Supply Chain Security"),
        ("12", "Security Monitoring & Logging"),
        ("13", "Awareness & Training"),
    ]
    col1_colors = [WS_NAVY, WS_BLUE_LT, WS_NAVY, WS_BLUE_LT, WS_NAVY, WS_BLUE_LT, WS_NAVY]
    col2_colors = [WS_BLUE_LT, WS_NAVY, WS_BLUE_LT, WS_NAVY, WS_BLUE_LT, WS_NAVY]
    row_h = Inches(0.42)
    start_y = Inches(1.33)

    for i, (num, name) in enumerate(cats[:7]):
        y = start_y + i * row_h
        rect(slide, Inches(0.45), y, Inches(0.38), Inches(0.34), fill_rgb=col1_colors[i])
        txbox(slide, num, Inches(0.45), y + Inches(0.02), Inches(0.38), Inches(0.30),
              font_size=Pt(7.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(slide, name, Inches(0.88), y + Inches(0.02), Inches(2.50), Inches(0.34),
              font_size=Pt(8.5), color=WS_DARK)

    for i, (num, name) in enumerate(cats[7:]):
        y = start_y + i * row_h
        rect(slide, Inches(3.50), y, Inches(0.38), Inches(0.34), fill_rgb=col2_colors[i])
        txbox(slide, num, Inches(3.50), y + Inches(0.02), Inches(0.38), Inches(0.30),
              font_size=Pt(7.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(slide, name, Inches(3.93), y + Inches(0.02), Inches(2.50), Inches(0.34),
              font_size=Pt(8.5), color=WS_DARK)

    # Classification output fields
    rect(slide, Inches(0.45), Inches(4.38), Inches(5.95), Inches(0.48),
         fill_rgb=RGBColor(0xEB, 0xF5, 0xFB), line_rgb=WS_BLUE_LT)
    txbox(slide, "Résultat : primary_category  ·  secondary_categories (max 2)  ·  confidence  ·  status",
          Inches(0.55), Inches(4.41), Inches(5.75), Inches(0.40),
          font_size=Pt(8.5), color=WS_NAVY)

    statuses = [
        ("validated",         "conf >= 0.85 et margin >= 0.15",       WS_GREEN),
        ("medium_confidence", "0.65 <= conf < 0.85",                   WS_YELLOW),
        ("multi_domain",      "2 categories proches (margin < 0.15)",  WS_ORANGE),
        ("low_confidence",    "conf < 0.65  ->  LLM fallback",         WS_RED),
    ]
    for i, (s, desc, c) in enumerate(statuses):
        y = Inches(4.95) + i * Inches(0.38)
        rect(slide, Inches(0.45), y, Inches(1.78), Inches(0.30), fill_rgb=c)
        txbox(slide, s, Inches(0.48), y + Inches(0.02), Inches(1.72), Inches(0.26),
              font_size=Pt(7.5), bold=True, color=WHITE)
        txbox(slide, desc, Inches(2.28), y + Inches(0.02), Inches(4.05), Inches(0.30),
              font_size=Pt(8), color=WS_MID)

    # ── RIGHT TOP: 4-step classification process ───────────────────────────
    rect(slide, Inches(6.65), Inches(0.85), Inches(6.30), Inches(3.30),
         fill_rgb=WHITE, line_rgb=WS_ORANGE)
    txbox(slide, "Processus de classification (4 etapes sequentielles)",
          Inches(6.75), Inches(0.90), Inches(6.10), Inches(0.35),
          font_size=Pt(11), bold=True, color=WS_DARK)

    steps = [
        ("1", "exact_match",
         "Correspondance mot-cle exacte  ->  conf = 1.0",            WS_GREEN),
        ("2", "override_rules",
         "Regles BUILTIN_RULES (deterministes)  ->  conf >= 0.98",   WS_BLUE_LT),
        ("3", "deterministic_score",
         "Score TF-IDF pondere  ->  conf variable, classement marge", WS_NAVY),
        ("4", "LLM fallback",
         "Si conf < 0.85 OU margin < 0.15  ->  appel Azure OpenAI",  WS_ORANGE),
    ]
    for i, (n, name, desc, color) in enumerate(steps):
        y = Inches(1.33) + i * Inches(0.64)
        rect(slide, Inches(6.75), y, Inches(0.42), Inches(0.52), fill_rgb=color)
        txbox(slide, n, Inches(6.75), y + Inches(0.08), Inches(0.42), Inches(0.36),
              font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, Inches(7.22), y, Inches(5.63), Inches(0.52),
             fill_rgb=RGBColor(0xF8, 0xF8, 0xF8), line_rgb=color)
        txbox(slide, name, Inches(7.32), y + Inches(0.02), Inches(5.45), Inches(0.22),
              font_size=Pt(9.5), bold=True, color=WS_DARK)
        txbox(slide, desc, Inches(7.32), y + Inches(0.26), Inches(5.45), Inches(0.22),
              font_size=Pt(8.5), color=WS_MID)

    # ── RIGHT BOTTOM: how ENISA is used in the pipeline ────────────────────
    rect(slide, Inches(6.65), Inches(4.27), Inches(6.30), Inches(2.63),
         fill_rgb=WHITE, line_rgb=WS_GREEN)
    txbox(slide, "Utilisation de l'ENISA dans le pipeline",
          Inches(6.75), Inches(4.32), Inches(6.10), Inches(0.35),
          font_size=Pt(11), bold=True, color=WS_GREEN)

    usages = [
        ("Pivot A <-> B",
         "Classification independante de chaque exigence\nRegroupe les paires de meme categorie ENISA"),
        ("Prior ENISA  (poids x0.03)",
         "primary==primary -> 1.0  |  primary in secondary -> 0.70\ncross-secondary -> 0.45  |  plafond 0.30 si conf faible"),
        ("Juge Final — groupement",
         "Decisions regroupees par source_category\nbatch_size=25 pour le prompt PROMPT_FINAL_JUDGE"),
        ("Mode soft_enisa",
         "Matching partiel inter-categories autorise\nPenalite reduite pour categories proches"),
    ]
    for i, (title, body) in enumerate(usages):
        row = i // 2
        col = i % 2
        x = Inches(6.75) + col * Inches(3.12)
        y = Inches(4.77) + row * Inches(1.05)
        bw = Inches(2.98)
        rect(slide, x, y, bw, Inches(0.95),
             fill_rgb=RGBColor(0xF0, 0xFB, 0xF4), line_rgb=WS_GREEN)
        txbox(slide, title, x + Inches(0.08), y + Inches(0.04), bw - Inches(0.14), Inches(0.26),
              font_size=Pt(9.5), bold=True, color=WS_GREEN)
        txbox(slide, body, x + Inches(0.08), y + Inches(0.30), bw - Inches(0.14), Inches(0.60),
              font_size=Pt(8), color=WS_MID)

    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), fill_rgb=WS_DARK)
    txbox(slide, "Regulatory Framework Mapper  |  Classification ENISA — Taxonomie Pivot",
          Inches(0.4), Inches(7.0), Inches(8), Inches(0.4),
          font_size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xAA))


def make_matching_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_LIGHT)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), fill_rgb=WS_BLUE_LT)
    rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill_rgb=WS_ORANGE)

    add_label(slide, "02", Inches(0.35), Inches(0.13), Inches(0.48), Inches(0.48),
              bg=WS_ORANGE, font_size=Pt(15), bold=True, align=PP_ALIGN.CENTER, rounded=True)
    txbox(slide, "Etape 2 — Scoring Multi-dimensionnel & Evaluation LLM par Paires",
          Inches(1.0), Inches(0.12), Inches(11.5), Inches(0.5),
          font_size=Pt(18), bold=True, color=WHITE)

    # ── LEFT: 5 scoring dimensions ──────────────────────────────────────────
    rect(slide, Inches(0.35), Inches(0.85), Inches(7.20), Inches(6.05),
         fill_rgb=WHITE, line_rgb=WS_BLUE_LT)
    txbox(slide, "Scoring multi-dimensionnel — 5 dimensions fusionnees par RRF (K=60)",
          Inches(0.45), Inches(0.90), Inches(7.00), Inches(0.35),
          font_size=Pt(11), bold=True, color=WS_NAVY)

    # D1: Semantic
    d1y = Inches(1.33)
    rect(slide, Inches(0.45), d1y, Inches(7.00), Inches(0.75),
         fill_rgb=RGBColor(0xEB, 0xF5, 0xFB), line_rgb=WS_BLUE_LT)
    txbox(slide, "D1  Semantique  x0.45",
          Inches(0.55), d1y + Inches(0.05), Inches(7.00), Inches(0.25),
          font_size=Pt(10), bold=True, color=WS_BLUE_LT)
    txbox(slide, "Cosinus(embed_source, embed_target) — champs ponderes :",
          Inches(0.55), d1y + Inches(0.30), Inches(7.00), Inches(0.20),
          font_size=Pt(8.5), color=WS_MID)
    txbox(slide, "object x3.0  |  action x2.5  |  control_type x1.8  |  evidence x1.5  |  actor/cond/deadline x1.0  |  obligation_type x0.8  |  domain x0.5",
          Inches(0.55), d1y + Inches(0.50), Inches(7.00), Inches(0.22),
          font_size=Pt(8), color=WS_DARK, italic=True)

    # D2: Structured
    d2y = Inches(2.16)
    rect(slide, Inches(0.45), d2y, Inches(7.00), Inches(0.75),
         fill_rgb=RGBColor(0xEB, 0xF5, 0xFB), line_rgb=WS_NAVY)
    txbox(slide, "D2  Structurel  x0.30",
          Inches(0.55), d2y + Inches(0.05), Inches(7.00), Inches(0.25),
          font_size=Pt(10), bold=True, color=WS_NAVY)
    txbox(slide, "Jaccard / overlap champ par champ — memes pondérations que D1 :",
          Inches(0.55), d2y + Inches(0.30), Inches(7.00), Inches(0.20),
          font_size=Pt(8.5), color=WS_MID)
    txbox(slide, "object x3.0  |  action x2.5  |  control_type x1.8  |  evidence x1.5  |  actor/cond/deadline x1.0  |  obligation_type x0.8  |  domain x0.5",
          Inches(0.55), d2y + Inches(0.50), Inches(7.00), Inches(0.22),
          font_size=Pt(8), color=WS_DARK, italic=True)

    # D3: Action/Object
    d3y = Inches(2.99)
    rect(slide, Inches(0.45), d3y, Inches(7.00), Inches(1.18),
         fill_rgb=RGBColor(0xFF, 0xF3, 0xEE), line_rgb=WS_ORANGE)
    txbox(slide, "D3  Action / Objet  x0.20",
          Inches(0.55), d3y + Inches(0.05), Inches(7.00), Inches(0.25),
          font_size=Pt(10), bold=True, color=WS_ORANGE)
    txbox(slide, "Jaccard(action U object U control_type)  +  bonus famille de controle +/-0.25",
          Inches(0.55), d3y + Inches(0.30), Inches(7.00), Inches(0.20),
          font_size=Pt(8.5), color=WS_MID)
    txbox(slide, "7 familles : access_control | awareness_training | incident | logging_monitoring | cryptography | backup_continuity | asset_inventory",
          Inches(0.55), d3y + Inches(0.50), Inches(7.00), Inches(0.22),
          font_size=Pt(8), color=WS_DARK, italic=True)
    rect(slide, Inches(0.55), d3y + Inches(0.76), Inches(6.85), Inches(0.34),
         fill_rgb=RGBColor(0xFF, 0xE0, 0xD0), line_rgb=WS_ORANGE)
    txbox(slide, "Gate AO : score < 0.10 -> cap final 30%   |   score < 0.25 -> cap final 50%   |   score >= 0.25 -> pas de cap",
          Inches(0.65), d3y + Inches(0.78), Inches(6.65), Inches(0.28),
          font_size=Pt(8.5), bold=True, color=WS_ORANGE)

    # D4: Control Type
    d4y = Inches(4.25)
    rect(slide, Inches(0.45), d4y, Inches(7.00), Inches(0.58),
         fill_rgb=RGBColor(0xF8, 0xF9, 0xFA), line_rgb=WS_MID)
    txbox(slide, "D4  Control Type  x0.05",
          Inches(0.55), d4y + Inches(0.05), Inches(7.00), Inches(0.25),
          font_size=Pt(10), bold=True, color=WS_MID)
    txbox(slide, "Jaccard pur sur control_type uniquement  |  Evite confusions : access_control <-> awareness_training",
          Inches(0.55), d4y + Inches(0.30), Inches(7.00), Inches(0.24),
          font_size=Pt(8.5), color=WS_MID)

    # D5: Prior ENISA
    d5y = Inches(4.91)
    rect(slide, Inches(0.45), d5y, Inches(7.00), Inches(0.75),
         fill_rgb=RGBColor(0xF0, 0xFB, 0xF4), line_rgb=WS_GREEN)
    txbox(slide, "D5  Prior ENISA  x0.03",
          Inches(0.55), d5y + Inches(0.05), Inches(7.00), Inches(0.25),
          font_size=Pt(10), bold=True, color=WS_GREEN)
    txbox(slide, "primary==primary -> 1.0   |   primary in secondary -> 0.70   |   cross-secondary -> 0.45",
          Inches(0.55), d5y + Inches(0.30), Inches(7.00), Inches(0.22),
          font_size=Pt(8.5), color=WS_MID)
    txbox(slide, "Boost uniquement — jamais penalite  |  Plafonne a 0.30 si confidence categorie faible",
          Inches(0.55), d5y + Inches(0.52), Inches(7.00), Inches(0.20),
          font_size=Pt(8), color=WS_DARK, italic=True)

    # RRF + obvious-gap shortcut note
    rect(slide, Inches(0.45), Inches(5.74), Inches(7.00), Inches(0.38),
         fill_rgb=RGBColor(0xEB, 0xF5, 0xFB), line_rgb=WS_BLUE_LT)
    txbox(slide, "Fusion RRF (K=60)  |  Gate MIN_SCORE=0.04  |  Gap evident : combined<0.10 + sem<0.20 + AO<0.15 -> pas d'appel LLM",
          Inches(0.55), Inches(5.76), Inches(6.80), Inches(0.32),
          font_size=Pt(8), color=WS_NAVY)

    # ── RIGHT: Coverage table (0 to 100 by 10%) ────────────────────────────
    rect(slide, Inches(7.70), Inches(0.85), Inches(5.45), Inches(6.05),
         fill_rgb=WHITE, line_rgb=WS_ORANGE)
    txbox(slide, "Niveaux de couverture — echelle 0 -> 100% par pas de 10%",
          Inches(7.80), Inches(0.90), Inches(5.25), Inches(0.35),
          font_size=Pt(11), bold=True, color=WS_DARK)

    cov_rows = [
        (100, "direct_full_coverage",          WS_GREEN),
        ( 90, "couverture quasi-complete",      RGBColor(0x27, 0xAE, 0x60)),
        ( 80, "mostly_covered (fort)",          RGBColor(0x52, 0xBE, 0x80)),
        ( 70, "mostly_covered",                 RGBColor(0x82, 0xE0, 0xAA)),
        ( 60, "partiellement couvert (bon)",    WS_YELLOW),
        ( 50, "partial / impl. detail gap",     WS_ACCENT),
        ( 40, "couverture faible",              RGBColor(0xF0, 0xB2, 0x7A)),
        ( 30, "indirect_support (possible)",    WS_RED),
        ( 20, "indirect_support",               RGBColor(0xE7, 0x4C, 0x3C)),
        ( 10, "couverture minimale",            RGBColor(0xC0, 0x39, 0x2B)),
        (  0, "true_gap / conflict",            RGBColor(0x44, 0x44, 0x44)),
    ]
    rh = Inches(0.44)
    for i, (cov, label, color) in enumerate(cov_rows):
        y = Inches(1.33) + i * rh
        rect(slide, Inches(7.80), y, Inches(0.72), rh - Inches(0.04), fill_rgb=color)
        txbox(slide, f"{cov}%", Inches(7.80), y + Inches(0.04), Inches(0.72), rh - Inches(0.10),
              font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, Inches(8.57), y, Inches(4.48), rh - Inches(0.04),
             fill_rgb=RGBColor(0xF9, 0xF9, 0xF9), line_rgb=color)
        txbox(slide, label, Inches(8.67), y + Inches(0.04), Inches(4.30), rh - Inches(0.10),
              font_size=Pt(8.5), color=WS_DARK)

    rect(slide, Inches(7.80), Inches(6.18), Inches(5.25), Inches(0.38),
         fill_rgb=RGBColor(0xFF, 0xF8, 0xE1), line_rgb=WS_YELLOW)
    txbox(slide, "Cache : mapping_decisions_cache.jsonl  |  Self-review optionnel (conf < 0.65)",
          Inches(7.90), Inches(6.20), Inches(5.05), Inches(0.32),
          font_size=Pt(8), color=RGBColor(0x7D, 0x6A, 0x00))

    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), fill_rgb=WS_DARK)
    txbox(slide, "Regulatory Framework Mapper  |  Etape 2 : Scoring Multi-dimensionnel",
          Inches(0.4), Inches(7.0), Inches(8), Inches(0.4),
          font_size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xAA))


def make_judge_synthesis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_LIGHT)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), fill_rgb=WS_ORANGE)
    rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill_rgb=WS_ORANGE)

    txbox(slide, "Etapes 3 & 4 — Juge Final  |  Synthese Parent  |  Post-traitement Automatique",
          Inches(0.35), Inches(0.12), Inches(12.5), Inches(0.5),
          font_size=Pt(18), bold=True, color=WHITE)

    # ── Stage 3: Final Judge (LEFT) ─────────────────────────────────────────
    rect(slide, Inches(0.35), Inches(0.85), Inches(6.15), Inches(4.25),
         fill_rgb=WHITE, line_rgb=WS_ORANGE)
    add_label(slide, "03", Inches(0.45), Inches(0.95), Inches(0.55), Inches(0.55),
              bg=WS_ORANGE, font_size=Pt(15), bold=True, align=PP_ALIGN.CENTER, rounded=True)
    txbox(slide, "Juge Final  (final_judge.py)",
          Inches(1.10), Inches(0.98), Inches(5.2), Inches(0.50),
          font_size=Pt(13), bold=True, color=WS_DARK)

    judge_items = [
        ("Filtre d'ambiguite",
         "Revise uniquement si : conf < 0.80  OU  coverage_level <= 50\n-> Passe directement si conf >= 0.80 ET coverage > 50"),
        ("Groupement & batch LLM",
         "Regroupe par source_category (ENISA)\nbatch_size = 25 decisions  |  Prompt : PROMPT_FINAL_JUDGE\nAppel Azure OpenAI (modele Judge)"),
        ("Corrections appliquees",
         "relation_type | coverage_level | confidence | match_type\ngap_items | selected_candidate_ids | target_requirements\n-> suivi de sanitize_decision() sur chaque correction"),
    ]
    for i, (title, body) in enumerate(judge_items):
        y = Inches(1.62) + i * Inches(1.12)
        rect(slide, Inches(0.45), y, Inches(5.95), Inches(1.02),
             fill_rgb=RGBColor(0xFF, 0xF3, 0xEE), line_rgb=WS_ORANGE)
        txbox(slide, title, Inches(0.55), y + Inches(0.04), Inches(5.75), Inches(0.26),
              font_size=Pt(9.5), bold=True, color=WS_ORANGE)
        txbox(slide, body, Inches(0.55), y + Inches(0.32), Inches(5.75), Inches(0.66),
              font_size=Pt(8.5), color=WS_MID)

    # ── Stage 4: Parent Synthesis (RIGHT) ──────────────────────────────────
    rect(slide, Inches(6.65), Inches(0.85), Inches(6.30), Inches(4.25),
         fill_rgb=WHITE, line_rgb=RGBColor(0x8E, 0x44, 0xAD))
    add_label(slide, "04", Inches(6.75), Inches(0.95), Inches(0.55), Inches(0.55),
              bg=RGBColor(0x8E, 0x44, 0xAD), font_size=Pt(15), bold=True,
              align=PP_ALIGN.CENTER, rounded=True)
    txbox(slide, "Synthese Parent  (parent_gap_synthesis.py)",
          Inches(7.40), Inches(0.98), Inches(5.40), Inches(0.50),
          font_size=Pt(13), bold=True, color=WS_DARK)

    synth_items = [
        ("Agregation Atomes -> Parent",
         "Combine les decisions atomiques par exigence parente\n[A1:full+100] + [A2:partial+50] + [A3:gap+0]\n-> coverage_level parent agrege par formule"),
        ("Types de lacune parent",
         "none | implementation_detail_gap | partial_gap\nindirect_support_gap | true_gap | conflict_gap\nDedoublonnage des gap_items par dimension"),
        ("Synthese LLM optionnelle",
         "Prompt PROMPT_PARENT_GAP_SYNTHESIS\nTexte de lacune redige par LLM\nCache : parent_gap_synthesis_cache.jsonl"),
    ]
    for i, (title, body) in enumerate(synth_items):
        y = Inches(1.62) + i * Inches(1.12)
        rect(slide, Inches(6.75), y, Inches(6.10), Inches(1.02),
             fill_rgb=RGBColor(0xF5, 0xEE, 0xF8), line_rgb=RGBColor(0x8E, 0x44, 0xAD))
        txbox(slide, title, Inches(6.85), y + Inches(0.04), Inches(5.90), Inches(0.26),
              font_size=Pt(9.5), bold=True, color=RGBColor(0x8E, 0x44, 0xAD))
        txbox(slide, body, Inches(6.85), y + Inches(0.32), Inches(5.90), Inches(0.66),
              font_size=Pt(8.5), color=WS_MID)

    arrow_right(slide, Inches(6.50), Inches(3.00), Inches(0.15), color=WS_ORANGE, width=Pt(2.5))

    # ── BOTTOM: Post-processing chain (sanitize_decision, 6 steps) ──────────
    rect(slide, Inches(0.35), Inches(5.22), Inches(12.65), Inches(1.58),
         fill_rgb=WHITE, line_rgb=WS_BLUE_LT)
    txbox(slide, "Post-traitement automatique — sanitize_decision() — applique apres chaque decision LLM et correction Juge",
          Inches(0.45), Inches(5.26), Inches(12.40), Inches(0.28),
          font_size=Pt(10), bold=True, color=WS_NAVY)

    post_steps = [
        ("Snap",      "Arrondi au\n10% le plus\nproche",               WS_BLUE_LT),
        ("Gate AO",   "AO < 0.10\n-> cap 30%\nAO < 0.25 -> 50%",      WS_ORANGE),
        ("Rescue gap","combined>=0.35\nou sem>=0.48\n-> indirect 20%", WS_GREEN),
        ("Upgrade",   "Tier1: comb>=0.58\n+sem>=0.76 -> 80%\nTier2: +sem>=0.84 -> 100%", WS_NAVY),
        ("Floor",     "comb>=0.60+cov<=30\n-> floor 50%\nsem>=0.80+AO>=0.20 -> 80%",     WS_MID),
        ("Downgrade", "cov<=20 : comb<0.45\nou sem<0.65\ncov=30 : les deux",              WS_RED),
    ]
    box_w = Inches(2.05)
    for i, (num, body, color) in enumerate(post_steps):
        x = Inches(0.45) + i * Inches(2.10)
        y = Inches(5.60)
        rect(slide, x, y, box_w, Inches(1.12),
             fill_rgb=RGBColor(0xF8, 0xF9, 0xFA), line_rgb=color)
        txbox(slide, f"{i+1}. {num}", x + Inches(0.06), y + Inches(0.04), box_w - Inches(0.10), Inches(0.26),
              font_size=Pt(9), bold=True, color=color)
        txbox(slide, body, x + Inches(0.06), y + Inches(0.30), box_w - Inches(0.10), Inches(0.78),
              font_size=Pt(7.5), color=WS_MID)
        if i < 5:
            arrow_right(slide, x + box_w, y + Inches(0.52), Inches(0.05), color=color, width=Pt(1.5))

    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), fill_rgb=WS_DARK)
    txbox(slide, "Regulatory Framework Mapper  |  Etapes 3 & 4 : Juge Final & Synthese Parent",
          Inches(0.4), Inches(7.0), Inches(9), Inches(0.4),
          font_size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xAA))


def make_output_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_LIGHT)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), fill_rgb=WS_GREEN)
    rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill_rgb=WS_ORANGE)

    add_label(slide, "05", Inches(0.35), Inches(0.13), Inches(0.48), Inches(0.48),
              bg=WS_ORANGE, font_size=Pt(15), bold=True, align=PP_ALIGN.CENTER, rounded=True)
    txbox(slide, "Étape 5 — Workbook Excel de sortie  (output_writer.py)",
          Inches(1.0), Inches(0.12), Inches(11), Inches(0.5),
          font_size=Pt(20), bold=True, color=WHITE)

    tabs = [
        ("README",             "Méthodologie, légende scoring,\ncatégories ENISA expliquées",       WS_MID,      "📖"),
        ("Dashboard",          "KPIs exécutifs + graphiques\nCoverage%, match types, gaps",         WS_BLUE_LT,  "📊"),
        ("Fr1 → Fr2\n(parent)", "Source · Cible · Coverage%\nLacune · Priorité review",            WS_NAVY,     "🗺"),
        ("Fr2 → Fr1\n(parent)", "Mapping inverse\n(si bidirectionnel)",                            WS_NAVY,     "🗺"),
        ("Coverage/Cat",       "Stats par catégorie ENISA\nTaux couverture agrégés",               WS_GREEN,    "📈"),
        ("Cat. Quality",       "Diagnostics catégories\nPiste qualité harmonisation",              WS_YELLOW,   "🔍"),
        ("Atomic Fr1→Fr2",     "Scores dim. · Rationale LLM\nCandidats diagnostics · Items gaps",  WS_ORANGE,   "⚛"),
        ("Atomic Fr2→Fr1",     "Piste d'évidence atomique\n(si bidirectionnel)",                   WS_ORANGE,   "⚛"),
    ]

    tw = Inches(1.5)
    th = Inches(2.5)
    tgap = Inches(0.12)
    sx = Inches(0.38)
    sy = Inches(0.95)

    for i, (name, desc, color, icon) in enumerate(tabs):
        x = sx + i * (tw + tgap)
        # tab card
        rounded_rect(slide, x, sy, tw, th, fill_rgb=WHITE, line_rgb=color)
        # colored top
        rounded_rect(slide, x, sy, tw, Inches(0.72), fill_rgb=color, line_rgb=None)
        txbox(slide, icon + " " + name, x + Inches(0.07), sy + Inches(0.08),
              tw - Inches(0.1), Inches(0.62),
              font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(slide, desc, x + Inches(0.07), sy + Inches(0.82),
              tw - Inches(0.1), Inches(1.55),
              font_size=Pt(8), color=WS_MID, align=PP_ALIGN.LEFT)

    # Column detail for parent sheet
    rect(slide, Inches(0.38), Inches(3.65), Inches(12.55), Inches(1.6),
         fill_rgb=WHITE, line_rgb=WS_NAVY)
    txbox(slide, "Colonnes de l'onglet parent (Fr1 → Fr2)",
          Inches(0.5), Inches(3.7), Inches(6), Inches(0.32),
          font_size=Pt(11), bold=True, color=WS_NAVY)

    cols_parent = [
        "Source Regulation", "Source Control ID", "Source Requirement", "ENISA Category",
        "Target Regulation", "Target Control ID(s)", "Target Requirement(s)",
        "Coverage Relationship", "Coverage Level (0-100)", "Gap", "Gap Détaillé", "Review Priority"
    ]
    for i, col in enumerate(cols_parent):
        row = i // 6
        col_i = i % 6
        cw = Inches(2.0)
        ch = Inches(0.32)
        cx = Inches(0.5) + col_i * (cw + Inches(0.05))
        cy = Inches(4.1) + row * (ch + Inches(0.05))
        rounded_rect(slide, cx, cy, cw, ch,
                     fill_rgb=RGBColor(0xEB, 0xF5, 0xFB), line_rgb=WS_BLUE_LT)
        txbox(slide, col, cx + Inches(0.05), cy + Inches(0.03),
              cw - Inches(0.08), ch - Inches(0.05),
              font_size=Pt(7.8), color=WS_NAVY, align=PP_ALIGN.CENTER)

    # Formatting note
    rect(slide, Inches(0.38), Inches(5.42), Inches(12.55), Inches(0.42),
         fill_rgb=RGBColor(0xEB, 0xF9, 0xF1), line_rgb=WS_GREEN)
    txbox(slide, "Mise en forme : color scales coverage% · conditional formatting · frozen headers · largeurs auto  "
                 "| Fichier : output/mapping_{FrA}_{FrB}_{date}_{time}.xlsx",
          Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.35),
          font_size=Pt(8.5), color=RGBColor(0x1A, 0x5C, 0x32))

    # Stage 6 compact
    rect(slide, Inches(0.38), Inches(5.98), Inches(12.55), Inches(0.72),
         fill_rgb=WHITE, line_rgb=WS_MID)
    add_label(slide, "06", Inches(0.48), Inches(6.05), Inches(0.48), Inches(0.48),
              bg=WS_MID, font_size=Pt(13), bold=True, align=PP_ALIGN.CENTER, rounded=True)
    txbox(slide, "Rapport d'analyse  (log_analyzer.py)",
          Inches(1.05), Inches(6.05), Inches(4), Inches(0.28),
          font_size=Pt(10), bold=True, color=WS_MID)
    txbox(slide, "logs/run_{...}.jsonl  →  reports/log_analysis_{...}.md   |   "
                 "Cache hit rate · Corrections juge · Distribution scores · Timing par étape · Résumé erreurs",
          Inches(1.05), Inches(6.37), Inches(11.7), Inches(0.28),
          font_size=Pt(8.5), color=WS_MID)

    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), fill_rgb=WS_DARK)
    txbox(slide, "Regulatory Framework Mapper  |  Étape 5 : Workbook Excel",
          Inches(0.4), Inches(7.0), Inches(8), Inches(0.4),
          font_size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xAA))


def make_transversal_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_LIGHT)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), fill_rgb=WS_DARK)
    rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill_rgb=WS_ORANGE)

    txbox(slide, "Composants transversaux",
          Inches(0.5), Inches(0.12), Inches(10), Inches(0.5),
          font_size=Pt(20), bold=True, color=WHITE)

    # Azure OpenAI block
    rect(slide, Inches(0.35), Inches(0.9), Inches(8.2), Inches(2.7),
         fill_rgb=WHITE, line_rgb=WS_BLUE_LT)
    add_label(slide, "Azure OpenAI", Inches(0.45), Inches(0.95), Inches(3.0), Inches(0.45),
              bg=WS_BLUE_LT, font_size=Pt(12), bold=True, align=PP_ALIGN.CENTER, rounded=True)

    deployments = [
        ("Text (génération)",      "gpt-4.1-nano / gpt-5.4-nano\nTempérature : 0.1  |  JSON mode",               WS_BLUE_LT),
        ("Judge (révision)",        "gpt-5.4-nano ou reasoning model\nAppels judge_json()",                       WS_NAVY),
        ("Embeddings",             "text-embedding-3-large  |  512 dims\nBatches 64  |  cosine similarity",       WS_GREEN),
    ]
    for i, (name, desc, color) in enumerate(deployments):
        x = Inches(0.45) + i * Inches(2.6)
        rect(slide, x, Inches(1.55), Inches(2.45), Inches(1.8),
             fill_rgb=RGBColor(0xF0, 0xF8, 0xFF), line_rgb=color)
        txbox(slide, name, x + Inches(0.1), Inches(1.6), Inches(2.3), Inches(0.28),
              font_size=Pt(9.5), bold=True, color=color)
        txbox(slide, desc, x + Inches(0.1), Inches(1.92), Inches(2.3), Inches(1.3),
              font_size=Pt(8.5), color=WS_MID)

    txbox(slide, "MAX_CONCURRENT_LLM_CALLS=6  (ThreadPoolExecutor)  |  generate_json()  ·  judge_json()  ·  embed_texts()",
          Inches(0.5), Inches(3.35), Inches(7.9), Inches(0.22),
          font_size=Pt(8), color=WS_MID, italic=True)

    # Cache block
    rect(slide, Inches(0.35), Inches(3.75), Inches(8.2), Inches(2.2),
         fill_rgb=WHITE, line_rgb=WS_YELLOW)
    add_label(slide, "Système de Cache  (cache.py)", Inches(0.45), Inches(3.8), Inches(3.5), Inches(0.45),
              bg=WS_YELLOW, txt_color=WS_DARK, font_size=Pt(12), bold=True,
              align=PP_ALIGN.CENTER, rounded=True)

    caches = [
        ("Cache Framework",
         "docs/cache/{framework_key}/\n• category_harmonization.json\n• atomized_requirements.json\n• atomic_requirements_with_fields.json\n• Validation SHA256 fichier source",
         WS_YELLOW),
        ("Cache Décisions",
         "mapping_decisions_cache.jsonl\nparent_gap_synthesis_cache.jsonl\n• Format JSONL  (thread-safe)\n• Évite appels LLM redondants\n• Clé = hash(source, candidat)",
         WS_ORANGE),
    ]
    for i, (name, desc, color) in enumerate(caches):
        x = Inches(0.45) + i * Inches(3.9)
        rect(slide, x, Inches(4.4), Inches(3.7), Inches(1.35),
             fill_rgb=RGBColor(0xFF, 0xFC, 0xF0), line_rgb=color)
        txbox(slide, name, x + Inches(0.1), Inches(4.45), Inches(3.5), Inches(0.28),
              font_size=Pt(9.5), bold=True, color=color)
        txbox(slide, desc, x + Inches(0.1), Inches(4.75), Inches(3.5), Inches(0.92),
              font_size=Pt(8.2), color=WS_MID)

    # RAG block
    rect(slide, Inches(8.75), Inches(0.9), Inches(4.2), Inches(5.05),
         fill_rgb=WHITE, line_rgb=RGBColor(0x1A, 0x8C, 0x6E))
    add_label(slide, "Guideline RAG  (optionnel)", Inches(8.85), Inches(0.95), Inches(4.0), Inches(0.45),
              bg=RGBColor(0x1A, 0x8C, 0x6E), font_size=Pt(12), bold=True,
              align=PP_ALIGN.CENTER, rounded=True)

    rag_steps = [
        "📂  Sources : pdf · txt · md · docx",
        "✂️  Chunking : 400 cars, overlap 80",
        "🔢  Embeddings par chunk",
        "🔍  Retrieval : Top-K passages",
        "💉  Injection dans PROMPT_PAIRWISE",
        "💾  Cache invalidé sur hash fichier",
    ]
    for i, step in enumerate(rag_steps):
        rect(slide, Inches(8.88), Inches(1.58) + i * Inches(0.63), Inches(3.9), Inches(0.52),
             fill_rgb=RGBColor(0xE8, 0xF8, 0xF4), line_rgb=RGBColor(0x1A, 0x8C, 0x6E))
        txbox(slide, step, Inches(9.0), Inches(1.62) + i * Inches(0.63), Inches(3.7), Inches(0.42),
              font_size=Pt(9), color=WS_DARK)

    # Validation block
    rect(slide, Inches(0.35), Inches(6.05), Inches(8.2), Inches(0.62),
         fill_rgb=RGBColor(0xFD, 0xED, 0xEC), line_rgb=WS_RED)
    txbox(slide, "Validation  (validation.py)  :  validate_non_empty_framework()  ·  validate_mapping_targets()  ·  "
                 "framework_atom_fingerprint()  |  Politiques : raise · drop · sanitize · convert_to_gap",
          Inches(0.5), Inches(6.08), Inches(8.0), Inches(0.55),
          font_size=Pt(8.5), color=WS_RED)

    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), fill_rgb=WS_DARK)
    txbox(slide, "Regulatory Framework Mapper  |  Composants Transversaux",
          Inches(0.4), Inches(7.0), Inches(8), Inches(0.4),
          font_size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xAA))


def make_config_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_LIGHT)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), fill_rgb=WS_MID)
    rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill_rgb=WS_ORANGE)

    txbox(slide, "Configuration & Paramétrage  (.env — 160+ paramètres)",
          Inches(0.5), Inches(0.12), Inches(11), Inches(0.5),
          font_size=Pt(20), bold=True, color=WHITE)

    groups = [
        ("ENTRÉES", WS_NAVY, [
            ("FRAMEWORK_A_FILE", "data/Test1_Belgique.xlsx"),
            ("FRAMEWORK_B_FILE", "data/Test2_France.xlsx"),
            ("A_ID_COLUMN / A_REQUIREMENT_COLUMN", "ID · Requirement"),
            ("MAX_REQUIREMENTS_PER_FRAMEWORK", "500"),
        ]),
        ("HARMONISATION ENISA", WS_BLUE_LT, [
            ("ENABLE_CATEGORY_HARMONIZATION", "true"),
            ("CATEGORY_HARMONIZATION_USE_LLM", "true"),
            ("MATCH_SCOPE", "soft_enisa (recommandé)"),
            ("STRONG/MEDIUM confidence", "0.85 / 0.60"),
        ]),
        ("POIDS SCORING", WS_ORANGE, [
            ("WEIGHT_SEMANTIC", "0.45"),
            ("WEIGHT_STRUCTURED", "0.30"),
            ("WEIGHT_ACTION_OBJECT", "0.20"),
            ("WEIGHT_CONTROL_TYPE", "0.05"),
            ("WEIGHT_CATEGORY_PRIOR", "0.03"),
        ]),
        ("ÉVALUATION LLM", RGBColor(0x8E, 0x44, 0xAD), [
            ("USE_LLM_PAIRWISE_EVALUATION", "true"),
            ("LLM_TOP_K_CANDIDATES", "6"),
            ("LLM_CONFIDENCE_THRESHOLD", "0.65"),
            ("LLM_SELF_REVIEW", "false"),
            ("MAX_CONCURRENT_LLM_CALLS", "6"),
        ]),
        ("JUGE FINAL", WS_RED, [
            ("RUN_FINAL_LLM_JUDGE", "true"),
            ("FINAL_JUDGE_ONLY_AMBIGUOUS", "true"),
            ("FINAL_JUDGE_BATCH_SIZE", "25"),
            ("FINAL_JUDGE_CONFIDENCE_THRESHOLD", "0.80"),
        ]),
        ("AZURE OPENAI", WS_GREEN, [
            ("TEXT_DEPLOYMENT", "gpt-4.1-nano"),
            ("JUDGE_DEPLOYMENT", "gpt-5.4-nano"),
            ("EMBEDDING_DEPLOYMENT", "text-emb-3-large"),
            ("EMBEDDING_DIMENSIONS", "512"),
            ("TEMPERATURE", "0.1"),
        ]),
    ]

    gw = Inches(2.1)
    gh = Inches(4.8)
    gx0 = Inches(0.35)
    gy0 = Inches(0.9)
    ggap = Inches(0.12)

    for i, (title, color, params) in enumerate(groups):
        x = gx0 + i * (gw + ggap)
        rounded_rect(slide, x, gy0, gw, gh, fill_rgb=WHITE, line_rgb=color)
        rounded_rect(slide, x, gy0, gw, Inches(0.45), fill_rgb=color, line_rgb=None)
        txbox(slide, title, x + Inches(0.06), gy0 + Inches(0.06),
              gw - Inches(0.1), Inches(0.35),
              font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        for j, (key, val) in enumerate(params):
            py = gy0 + Inches(0.55) + j * Inches(0.82)
            rect(slide, x + Inches(0.08), py, gw - Inches(0.14), Inches(0.72),
                 fill_rgb=RGBColor(0xF8, 0xF8, 0xF8), line_rgb=RGBColor(0xCC, 0xCC, 0xCC))
            txbox(slide, key, x + Inches(0.13), py + Inches(0.04),
                  gw - Inches(0.2), Inches(0.3),
                  font_size=Pt(7.5), bold=True, color=WS_DARK)
            add_label(slide, val, x + Inches(0.13), py + Inches(0.36),
                      gw - Inches(0.22), Inches(0.28),
                      bg=color, txt_color=WHITE, font_size=Pt(7.5),
                      bold=False, align=PP_ALIGN.CENTER, rounded=True)

    # quality vs cost note
    rect(slide, Inches(0.35), Inches(5.9), Inches(12.55), Inches(0.8),
         fill_rgb=RGBColor(0xEB, 0xF5, 0xFB), line_rgb=WS_BLUE_LT)
    txbox(slide, "Mode économique  :  gpt-4.1-nano · TOP_K=5 · FINAL_JUDGE_ONLY_AMBIGUOUS=true · MAX_CONCURRENT=2\n"
                 "Mode qualité     :  gpt-4.1-mini · JUDGE=gpt-5.4-nano · LLM_SELF_REVIEW=true · FINAL_JUDGE_ONLY_AMBIGUOUS=false",
          Inches(0.5), Inches(5.93), Inches(12.3), Inches(0.72),
          font_size=Pt(9), color=WS_NAVY)

    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), fill_rgb=WS_DARK)
    txbox(slide, "Regulatory Framework Mapper  |  Configuration & Paramétrage",
          Inches(0.4), Inches(7.0), Inches(8), Inches(0.4),
          font_size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xAA))


def make_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_DARK)
    rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill_rgb=WS_ORANGE)

    rect(slide, Inches(0), Inches(4.8), SLIDE_W, Inches(2.7), fill_rgb=WS_NAVY)

    txbox(slide, "Points clés de l'architecture",
          Inches(0.7), Inches(0.5), Inches(11), Inches(0.55),
          font_size=Pt(24), bold=True, color=WHITE)

    innovations = [
        ("Scoring hybride",          "5 dimensions combinées par RRF\nSémantique + champs structurés + Action/Objet gate"),
        ("Jugement LLM progressif",  "Pairwise → Self-review → Juge final\nConfidence-weighted, batch par catégorie"),
        ("Caching agressif",         "SHA256 validation framework\nDécisions JSONL réutilisées entre runs"),
        ("Traçabilité totale",        "Atomes → Parent → Dashboard\nDimension-level evidence trail"),
        ("ENISA comme pivot",        "13 catégories comme langue commune\nAdvisory prior, jamais hard filter"),
        ("RAG guidelines",           "Contexte réglementaire injecté\nDans chaque évaluation pairwise"),
    ]

    for i, (title, desc) in enumerate(innovations):
        col = i % 3
        row = i // 3
        x = Inches(0.7) + col * Inches(4.1)
        y = Inches(1.3) + row * Inches(1.4)
        rect(slide, x, y, Inches(3.8), Inches(1.2),
             fill_rgb=RGBColor(0x22, 0x22, 0x40), line_rgb=WS_ORANGE)
        rect(slide, x, y, Inches(0.1), Inches(1.2), fill_rgb=WS_ORANGE)
        txbox(slide, title, x + Inches(0.2), y + Inches(0.08),
              Inches(3.5), Inches(0.35),
              font_size=Pt(11), bold=True, color=WS_ACCENT)
        txbox(slide, desc, x + Inches(0.2), y + Inches(0.45),
              Inches(3.5), Inches(0.65),
              font_size=Pt(9), color=RGBColor(0xCC, 0xCC, 0xCC))

    txbox(slide, "Wavestone  |  Regulatory Intelligence Platform  |  2026",
          Inches(0.7), Inches(6.95), Inches(9), Inches(0.4),
          font_size=Pt(10), color=RGBColor(0xAA, 0xAA, 0xAA))


# ── detailed pipeline slide ──────────────────────────────────────────────────

def make_detailed_pipeline_slide(prs):
    """Slide intermédiaire — 6 étapes avec sous-steps, modules, artefacts et params clés."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_rgb=WS_LIGHT)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.75), fill_rgb=WS_DARK)
    rect(slide, Inches(0), Inches(0), Inches(0.2), SLIDE_H, fill_rgb=WS_ORANGE)

    txbox(slide, "Pipeline détaillé — du Framework source à l'Excel de sortie",
          Inches(0.4), Inches(0.12), Inches(11), Inches(0.5),
          font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    stages = [
        {
            "num": "01", "title": "Prétraitement", "color": WS_NAVY,
            "module": "preprocessor.py  ·  embedder.py",
            "steps": [
                "Lecture onglet Excel (colonnes configurables)",
                "Mapping ID / Title / Requirement / Category",
                "Normalisation → 13 catégories ENISA",
                "  règles déterministes + LLM en fallback",
                "Atomisation LLM (obligations unitaires)",
                "Extraction 9 champs structurés par atome",
                "  actor · action · object · condition",
                "  deadline · evidence · obligation_type",
                "Embeddings text-embedding-3-large (512d)",
            ],
            "artifact": "atomic_requirements_with_fields.json",
            "param": "MAX_REQ • STRONG≥0.85 • BATCH=64",
        },
        {
            "num": "02", "title": "Mapping directionnel", "color": WS_BLUE_LT,
            "module": "matcher.py  ·  llm_evaluator.py",
            "steps": [
                "Score sémantique cosinus (×0.45)",
                "Score structuré — 5 champs (×0.30)",
                "Gate action/objet anti-FP (×0.20)",
                "Control-type gate (×0.05)",
                "Prior catégorie ENISA boost (×0.03)",
                "Fusion RRF K=60 → score hybride",
                "PROMPT_PAIRWISE_MATCH (Top-K candidats)",
                "Verdict : full / mostly / partial / gap",
                "Self-review optionnel (conf < seuil)",
            ],
            "artifact": "mapping_decisions_cache.jsonl",
            "param": "SEM=0.45 • TOP_K=6 • MIN=0.04",
        },
        {
            "num": "03", "title": "Juge final", "color": WS_ORANGE,
            "module": "final_judge.py",
            "steps": [
                "Groupement décisions / catégorie ENISA",
                "Batch 25 → PROMPT_FINAL_JUDGE",
                "Mode ambiguous_only (conf < 0.80)",
                "Correction FP : exact → partial",
                "Correction FN : true_gap → partial",
                "Scores dimensionnels affinés",
                "Journalisation des corrections",
                "",
                "",
            ],
            "artifact": "décisions révisées (en mémoire)",
            "param": "BATCH=25 • AMBIGUOUS_ONLY • CONF=0.80",
        },
        {
            "num": "04", "title": "Synthèse parent", "color": RGBColor(0x8E, 0x44, 0xAD),
            "module": "parent_gap_synthesis.py",
            "steps": [
                "Agrégation atomes → exigence parente",
                "  [A1:full]+[A2:partial] → partial_gap",
                "6 types de lacune parent classifiés :",
                "  none · partial_gap · true_gap",
                "  indirect_support · impl_detail · conflict",
                "Synthèse textuelle lacune via LLM",
                "Dédoublonnage items par dimension",
                "Cache synthèse JSONL (thread-safe)",
                "",
            ],
            "artifact": "parent_gap_synthesis_cache.jsonl",
            "param": "SYNTHESIS_LLM • DEDUP dim.",
        },
        {
            "num": "05", "title": "Workbook Excel", "color": WS_GREEN,
            "module": "output_writer.py",
            "steps": [
                "Onglet parent (12 colonnes structurées)",
                "  Coverage Level · Gap · Priority",
                "Onglets atomiques : scores dim. + rationale",
                "Dashboard KPIs + graphiques auto",
                "Coverage by Category (ENISA)",
                "Category Quality diagnostics",
                "Color scales + conditional formatting",
                "Frozen headers + largeurs adaptées",
                "Bidirectionnel : 8 onglets complets",
            ],
            "artifact": "mapping_{FrA}_{FrB}_{date}.xlsx",
            "param": "8 onglets • BIDIR optionnel",
        },
        {
            "num": "06", "title": "Rapport analyse", "color": WS_MID,
            "module": "log_analyzer.py",
            "steps": [
                "Log JSONL structuré par événement",
                "Cache hit rate calculé par étape",
                "Distribution scores et verdicts LLM",
                "Statistiques corrections juge final",
                "Timing par étape et par framework",
                "Résumé erreurs et avertissements",
                "Rapport Markdown auto-généré",
                "",
                "",
            ],
            "artifact": "reports/log_analysis_{...}.md",
            "param": "JSONL structuré • replay possible",
        },
    ]

    box_w = Inches(1.97)
    box_h = Inches(5.5)
    gap   = Inches(0.165)
    sx    = Inches(0.37)
    sy    = Inches(0.87)

    for i, stage in enumerate(stages):
        x     = sx + i * (box_w + gap)
        color = stage["color"]

        # Card body
        rounded_rect(slide, x, sy, box_w, box_h, fill_rgb=WHITE, line_rgb=color)

        # Colored header (rounded top + flat bottom overlap)
        rounded_rect(slide, x, sy, box_w, Inches(0.82), fill_rgb=color, line_rgb=None)
        rect(slide, x, sy + Inches(0.52), box_w, Inches(0.30), fill_rgb=color, line_rgb=None)

        # Stage number badge
        add_label(slide, stage["num"],
                  x + Inches(0.07), sy + Inches(0.1),
                  Inches(0.42), Inches(0.42),
                  bg=WHITE, txt_color=color, font_size=Pt(12), bold=True,
                  align=PP_ALIGN.CENTER, rounded=True)

        # Stage title
        txbox(slide, stage["title"],
              x + Inches(0.55), sy + Inches(0.12),
              box_w - Inches(0.60), Inches(0.62),
              font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

        # Module tag (light band)
        rect(slide, x, sy + Inches(0.82), box_w, Inches(0.28),
             fill_rgb=RGBColor(0xEC, 0xEC, 0xF2), line_rgb=None)
        txbox(slide, stage["module"],
              x + Inches(0.07), sy + Inches(0.84),
              box_w - Inches(0.12), Inches(0.22),
              font_size=Pt(7.2), bold=False, color=WS_NAVY, italic=True,
              align=PP_ALIGN.LEFT)

        # Steps
        steps_text = "\n".join(s for s in stage["steps"])
        txbox(slide, steps_text,
              x + Inches(0.08), sy + Inches(1.16),
              box_w - Inches(0.14), Inches(3.0),
              font_size=Pt(8.0), color=WS_DARK, align=PP_ALIGN.LEFT)

        # Divider
        rect(slide, x + Inches(0.1), sy + Inches(4.24),
             box_w - Inches(0.20), Inches(0.02),
             fill_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_rgb=None)

        # Artifact tag
        rect(slide, x + Inches(0.07), sy + Inches(4.30),
             box_w - Inches(0.12), Inches(0.52),
             fill_rgb=RGBColor(0xF0, 0xF4, 0xFF), line_rgb=color)
        txbox(slide, stage["artifact"],
              x + Inches(0.11), sy + Inches(4.33),
              box_w - Inches(0.20), Inches(0.44),
              font_size=Pt(7.0), color=WS_NAVY, italic=True, align=PP_ALIGN.LEFT)

        # Key-param badge
        rect(slide, x + Inches(0.07), sy + Inches(4.90),
             box_w - Inches(0.12), Inches(0.40),
             fill_rgb=color, line_rgb=None)
        txbox(slide, stage["param"],
              x + Inches(0.10), sy + Inches(4.93),
              box_w - Inches(0.18), Inches(0.34),
              font_size=Pt(7.0), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Arrow to next card
        if i < len(stages) - 1:
            ax = int(x + box_w + Inches(0.01))
            ay = int(sy + box_h / 2)
            arrow_right(slide, ax, ay, int(gap - Inches(0.02)),
                        color=WS_MID, width=Pt(2.5))

    # Bottom input / output strips
    rect(slide, Inches(0.25), Inches(6.48), Inches(4.8), Inches(0.34),
         fill_rgb=RGBColor(0xEB, 0xF5, 0xFB), line_rgb=WS_BLUE_LT)
    txbox(slide, "ENTREES :  Framework A + B (.xlsx)   ·   Guidelines RAG (pdf / txt / md)",
          Inches(0.35), Inches(6.50), Inches(4.6), Inches(0.28),
          font_size=Pt(7.5), color=WS_NAVY, bold=False)

    rect(slide, Inches(8.30), Inches(6.48), Inches(4.8), Inches(0.34),
         fill_rgb=RGBColor(0xEB, 0xF9, 0xF1), line_rgb=WS_GREEN)
    txbox(slide, "SORTIES :  mapping_{FrA}_{FrB}_{date}.xlsx   ·   log_analysis_{...}.md",
          Inches(8.40), Inches(6.50), Inches(4.6), Inches(0.28),
          font_size=Pt(7.5), color=RGBColor(0x1A, 0x5C, 0x32), bold=False)

    # Bottom bar
    rect(slide, Inches(0), Inches(6.95), SLIDE_W, Inches(0.55), fill_rgb=WS_DARK)
    txbox(slide, "Regulatory Framework Mapper  |  Pipeline Détaillé",
          Inches(0.4), Inches(7.0), Inches(8), Inches(0.4),
          font_size=Pt(9), color=RGBColor(0xAA, 0xAA, 0xAA))


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs)
    make_overview_slide(prs)
    make_detailed_pipeline_slide(prs)
    make_preprocessing_slide(prs)
    make_enisa_slide(prs)
    make_matching_slide(prs)
    make_judge_synthesis_slide(prs)
    make_output_slide(prs)
    make_transversal_slide(prs)
    make_config_slide(prs)
    make_closing_slide(prs)

    out = r"c:\Users\000023026\OneDrive - Wavestone\Desktop\mapper\Regulatory_Framework_Mapper_Architecture.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
